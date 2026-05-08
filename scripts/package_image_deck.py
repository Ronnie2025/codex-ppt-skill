#!/usr/bin/env python3
"""将整页生图封装为图片型 PPTX，并可选叠加真实 Logo。"""

from __future__ import annotations

import argparse
import math
import re
import shutil
import subprocess
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--images-dir", required=True, help="包含原始整页生图的目录。")
    parser.add_argument("--out-pptx", required=True, help="输出 PPTX 路径。")
    parser.add_argument("--final-dir", help="后处理后的最终 slide PNG 输出目录。")
    parser.add_argument("--contact-sheet", help="contact sheet 缩略图输出路径。")
    parser.add_argument("--slide-count", type=int, help="期望页数，用于检查缺页。")
    parser.add_argument("--width", type=int, default=1920, help="最终 slide 图片宽度。")
    parser.add_argument("--height", type=int, default=1080, help="最终 slide 图片高度。")
    parser.add_argument("--logo", help="每页叠加的 Logo，通常放在右上角。")
    parser.add_argument("--logo-width", type=int, default=160, help="每页 Logo 的像素宽度。")
    parser.add_argument("--logo-margin-x", type=int, default=70, help="每页 Logo 距右侧边距。")
    parser.add_argument("--logo-margin-y", type=int, default=45, help="每页 Logo 距顶部边距。")
    parser.add_argument("--mask-logo-zone", action="store_true", help="在右上角 Logo 后方绘制白色遮罩，覆盖预留框。")
    parser.add_argument(
        "--footer-logo",
        action="append",
        default=[],
        metavar="PATH:WIDTH",
        help="封面/尾页页脚 Logo，格式为 PATH:WIDTH，可重复传入。",
    )
    parser.add_argument("--footer-logo-slides", default="first,last", help="页脚 Logo 应用范围：first,last,all,none，可逗号分隔。")
    parser.add_argument("--export-pdf", action="store_true", help="如果本机有 soffice/libreoffice，则同时导出 PDF。")
    return parser.parse_args()


def natural_key(path: Path):
    return [int(s) if s.isdigit() else s.lower() for s in re.split(r"(\\d+)", path.name)]


def find_images(images_dir: Path) -> list[Path]:
    exts = {".png", ".jpg", ".jpeg", ".webp"}
    return sorted([p for p in images_dir.iterdir() if p.suffix.lower() in exts], key=natural_key)


def crop_to_size(img: Image.Image, width: int, height: int) -> Image.Image:
    img = img.convert("RGB")
    src_ratio = img.width / img.height
    dst_ratio = width / height
    if src_ratio > dst_ratio:
        new_h = height
        new_w = round(height * src_ratio)
    else:
        new_w = width
        new_h = round(width / src_ratio)
    resized = img.resize((new_w, new_h), Image.LANCZOS)
    left = (new_w - width) // 2
    top = (new_h - height) // 2
    return resized.crop((left, top, left + width, top + height))


def load_logo(spec: str) -> tuple[Image.Image, int]:
    if ":" in spec:
        path_s, width_s = spec.rsplit(":", 1)
        width = int(width_s)
    else:
        path_s = spec
        width = 180
    logo = Image.open(path_s).convert("RGBA")
    h = round(logo.height * (width / logo.width))
    return logo.resize((width, h), Image.LANCZOS), width


def overlay_logo(base: Image.Image, logo_path: str, width: int, mx: int, my: int, mask: bool) -> None:
    logo = Image.open(logo_path).convert("RGBA")
    h = round(logo.height * (width / logo.width))
    logo = logo.resize((width, h), Image.LANCZOS)
    x = base.width - width - mx
    y = my
    if mask:
        draw = ImageDraw.Draw(base)
        pad = 26
        draw.rounded_rectangle(
            [x - pad, y - pad, x + width + pad, y + h + pad],
            radius=12,
            fill=(255, 255, 255),
        )
    base.paste(logo, (x, y), logo)


def footer_slide_indices(count: int, mode: str) -> set[int]:
    parts = {p.strip().lower() for p in mode.split(",") if p.strip()}
    if "none" in parts:
        return set()
    if "all" in parts:
        return set(range(count))
    indices = set()
    if "first" in parts:
        indices.add(0)
    if "last" in parts:
        indices.add(count - 1)
    return indices


def overlay_footer_logos(base: Image.Image, specs: list[str]) -> None:
    if not specs:
        return
    logos = [load_logo(spec)[0] for spec in specs]
    x = 95
    y = base.height - 155
    mask_w = sum(l.width for l in logos) + 55 * max(0, len(logos) - 1) + 70
    mask_h = max(l.height for l in logos) + 46
    draw = ImageDraw.Draw(base)
    draw.rounded_rectangle([x - 25, y - 22, x - 25 + mask_w, y - 22 + mask_h], radius=10, fill=(255, 255, 255))
    for logo in logos:
        base.paste(logo, (x, y), logo)
        x += logo.width + 55


def make_contact_sheet(files: list[Path], out: Path) -> None:
    thumb_w, thumb_h, label_h = 384, 216, 34
    cols = 3
    rows = math.ceil(len(files) / cols)
    sheet = Image.new("RGB", (cols * thumb_w, rows * (thumb_h + label_h)), "white")
    draw = ImageDraw.Draw(sheet)
    for idx, file in enumerate(files):
        img = Image.open(file).convert("RGB").resize((thumb_w, thumb_h), Image.LANCZOS)
        col, row = idx % cols, idx // cols
        x, y = col * thumb_w, row * (thumb_h + label_h)
        sheet.paste(img, (x, y))
        draw.text((x + 10, y + thumb_h + 8), f"slide-{idx + 1:02d}", fill=(80, 80, 80))
    out.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(out, quality=92)


def build_pptx(files: list[Path], out_pptx: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for file in files:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(file), 0, 0, width=prs.slide_width, height=prs.slide_height)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)


def export_pdf(out_pptx: Path) -> None:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        print("跳过 PDF 导出：未找到 soffice/libreoffice。")
        return
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_pptx.parent), str(out_pptx)],
        check=True,
    )


def main() -> None:
    args = parse_args()
    images_dir = Path(args.images_dir)
    out_pptx = Path(args.out_pptx)
    final_dir = Path(args.final_dir) if args.final_dir else out_pptx.with_suffix("").parent / "final_slides"
    final_dir.mkdir(parents=True, exist_ok=True)

    images = find_images(images_dir)
    if args.slide_count and len(images) != args.slide_count:
        raise SystemExit(f"期望 {args.slide_count} 张图片，但在 {images_dir} 找到 {len(images)} 张。")
    if not images:
        raise SystemExit(f"目录中没有找到图片：{images_dir}")

    footer_indices = footer_slide_indices(len(images), args.footer_logo_slides)
    final_files = []
    for idx, img_path in enumerate(images):
        img = crop_to_size(Image.open(img_path), args.width, args.height)
        if args.logo:
            overlay_logo(img, args.logo, args.logo_width, args.logo_margin_x, args.logo_margin_y, args.mask_logo_zone)
        if idx in footer_indices:
            overlay_footer_logos(img, args.footer_logo)
        out = final_dir / f"slide-{idx + 1:02d}.png"
        img.save(out)
        final_files.append(out)

    build_pptx(final_files, out_pptx)
    if args.contact_sheet:
        make_contact_sheet(final_files, Path(args.contact_sheet))
    if args.export_pdf:
        export_pdf(out_pptx)
    print(f"已写入 {out_pptx}")


if __name__ == "__main__":
    main()
