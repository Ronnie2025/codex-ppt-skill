#!/usr/bin/env python3
"""Build an editable PPTX from semantic inventory and asset manifest files."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


EMU_PER_INCH = 914400
DEFAULT_SLIDE_W = 13.333333
DEFAULT_SLIDE_H = 7.5
VALID_ASSET_SOURCES = {"imagegen_asset", "api_generated_asset", "provided_asset"}
TEXT_WIDTH_FACTOR = 1.08
NO_WRAP_DEFAULT_ROLES = {"title", "subtitle", "header", "micro_label"}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--inventory", required=True, help="visual_inventory.json path.")
    parser.add_argument("--manifest", required=True, help="asset_manifest.json path.")
    parser.add_argument("--out-pptx", required=True, help="Output editable PPTX path.")
    parser.add_argument("--base-dir", help="Base directory for relative asset paths. Defaults to inventory parent.")
    parser.add_argument("--layout-rules", help="Optional layout_rules.json path for font and QA policy.")
    parser.add_argument("--slide-width", type=float, default=DEFAULT_SLIDE_W, help="Slide width in inches.")
    parser.add_argument("--slide-height", type=float, default=DEFAULT_SLIDE_H, help="Slide height in inches.")
    parser.add_argument("--report", help="Optional JSON build report path.")
    parser.add_argument("--min-body-font", type=float, default=5.8, help="QA warning floor for body text after fitting.")
    parser.add_argument("--min-title-font", type=float, default=8.0, help="QA error floor for title/header text after fitting.")
    parser.add_argument("--max-overflow-ratio", type=float, default=1.03, help="Allowed estimated text overflow ratio.")
    parser.add_argument("--collision-threshold", type=float, default=0.08, help="Minimum text overlap share to report as collision.")
    parser.add_argument("--fail-on-layout-qa", action="store_true", help="Exit non-zero when layout QA has errors.")
    return parser.parse_args()


def read_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def hex_to_rgb(value: str | None, default: str = "FFFFFF") -> RGBColor:
    value = (value or default).strip().lstrip("#")
    if len(value) != 6:
        value = default
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def rel_or_abs(base: Path, value: str) -> Path:
    path = Path(value)
    return path if path.is_absolute() else base / path


def px_box_to_inches(box: list[float], slide_size_px: list[float], slide_w: float, slide_h: float) -> tuple[float, float, float, float]:
    if len(box) != 4:
        raise ValueError(f"bbox must contain 4 numbers: {box}")
    sx = slide_w / float(slide_size_px[0])
    sy = slide_h / float(slide_size_px[1])
    return box[0] * sx, box[1] * sy, box[2] * sx, box[3] * sy


def as_emu(value_in: float) -> int:
    return int(round(value_in * EMU_PER_INCH))


def text_units(value: str) -> float:
    units = 0.0
    for char in value:
        code = ord(char)
        if char.isspace():
            units += 0.35
        elif code < 128:
            units += 0.55
        elif "\u3000" <= char <= "\u303f" or "\uff00" <= char <= "\uffef":
            units += 0.55
        else:
            units += 1.0
    return units


def wrap_line_by_units(value: str, max_units: float) -> list[str]:
    if not value:
        return [""]
    if max_units <= 1:
        return [value]
    lines: list[str] = []
    current = ""
    current_units = 0.0
    for char in value:
        unit = text_units(char)
        if current and current_units + unit > max_units:
            lines.append(current)
            current = char
            current_units = unit
        else:
            current += char
            current_units += unit
    if current:
        lines.append(current)
    return lines or [value]


def infer_text_role(item: dict[str, Any]) -> str:
    value = str(item.get("text_role") or item.get("role") or "").lower()
    if value:
        return value
    item_id = str(item.get("id") or "").lower()
    font_size = float(item.get("font_size", 14))
    if any(token in item_id for token in ("tag", "pill", "chip", "badge")):
        return "micro_label"
    if "title" in item_id or font_size >= 24:
        return "title"
    if "head" in item_id or "_h" in item_id or font_size >= 14:
        return "header"
    if "sub" in item_id:
        return "subtitle"
    if "body" in item_id or "_b" in item_id or font_size <= 9:
        return "body"
    return "label"


def role_min_font(role: str, item: dict[str, Any], args: argparse.Namespace) -> float:
    if item.get("min_font_size") is not None:
        return float(item["min_font_size"])
    if role == "micro_label":
        return 4.8
    if role == "header":
        return 7.0
    if role in {"title", "subtitle"}:
        return args.min_title_font
    return args.min_body_font


def line_step_pt(font_size: float, item: dict[str, Any]) -> float:
    if item.get("line_spacing_pt") is not None:
        return float(item["line_spacing_pt"])
    if item.get("line_spacing") is not None:
        raw = float(item["line_spacing"])
        return raw if raw > 3 else font_size * raw
    return font_size * float(item.get("line_height", 1.0))


def width_factor_for_role(role: str) -> float:
    if role == "title":
        return 1.32
    if role in {"subtitle", "header"}:
        return 1.16
    if role == "micro_label":
        return 1.10
    return TEXT_WIDTH_FACTOR


def text_margins_in(item: dict[str, Any], box: tuple[float, float, float, float] | None = None) -> tuple[float, float, float, float]:
    if box:
        _, _, w, h = box
        default_lr = 0.005 if w < 0.75 else 0.015 if w < 1.25 else 0.02
        default_tb = 0.0 if h < 0.13 else 0.004 if h < 0.2 else 0.01
    else:
        default_lr = 0.02
        default_tb = 0.01
    return (
        float(item.get("margin_left", default_lr)),
        float(item.get("margin_right", default_lr)),
        float(item.get("margin_top", default_tb)),
        float(item.get("margin_bottom", default_tb)),
    )


def estimate_text_layout(
    item: dict[str, Any],
    box: tuple[float, float, float, float],
    font_size: float,
) -> dict[str, Any]:
    x, y, w, h = box
    ml, mr, mt, mb = text_margins_in(item, box)
    inner_w_pt = max(1.0, (w - ml - mr) * 72.0)
    inner_h_pt = max(1.0, (h - mt - mb) * 72.0)
    role = infer_text_role(item)
    width_factor = width_factor_for_role(role)
    max_units = max(1.0, inner_w_pt / max(1.0, font_size * width_factor))
    word_wrap = bool(item.get("word_wrap", role not in NO_WRAP_DEFAULT_ROLES))
    visual_lines: list[str] = []
    for raw_line in str(item.get("text", "")).split("\n"):
        if word_wrap:
            visual_lines.extend(wrap_line_by_units(raw_line, max_units))
        else:
            visual_lines.append(raw_line)
    step = line_step_pt(font_size, item)
    estimated_h_pt = max(step, len(visual_lines) * step)
    longest_units = max((text_units(line) for line in visual_lines), default=0.0)
    estimated_w_pt = longest_units * font_size * width_factor
    overflow_h = estimated_h_pt / inner_h_pt
    overflow_w = estimated_w_pt / inner_w_pt
    vertical = str(item.get("vertical_anchor", "middle")).lower()
    content_h_in = estimated_h_pt / 72.0
    inner_x = x + ml
    inner_y = y + mt
    inner_w = max(0.001, w - ml - mr)
    inner_h = max(0.001, h - mt - mb)
    if vertical == "bottom":
        content_y = y + h - mb - content_h_in
    elif vertical == "top":
        content_y = inner_y
    else:
        content_y = y + (h - content_h_in) / 2.0
    return {
        "visual_line_count": len(visual_lines),
        "estimated_height_pt": round(estimated_h_pt, 3),
        "estimated_width_pt": round(estimated_w_pt, 3),
        "inner_width_pt": round(inner_w_pt, 3),
        "inner_height_pt": round(inner_h_pt, 3),
        "overflow_height_ratio": round(overflow_h, 4),
        "overflow_width_ratio": round(overflow_w, 4),
        "max_units_per_line": round(max_units, 3),
        "line_spacing_pt": round(step, 3),
        "inner_bbox_in": [round(inner_x, 4), round(inner_y, 4), round(inner_w, 4), round(inner_h, 4)],
        "estimated_text_bbox_in": [
            round(inner_x, 4),
            round(content_y, 4),
            round(min(inner_w, max(inner_w, estimated_w_pt / 72.0)), 4),
            round(content_h_in, 4),
        ],
        "wrapped_lines_preview": visual_lines[:8],
    }


def choose_effective_font(
    item: dict[str, Any],
    box: tuple[float, float, float, float],
    args: argparse.Namespace,
) -> tuple[float, dict[str, Any], str]:
    start = float(item.get("font_size", 14))
    role = infer_text_role(item)
    min_size = role_min_font(role, item, args)
    fit_mode = str(item.get("fit_mode", "shrink")).lower()
    if fit_mode not in {"shrink", "none"}:
        fit_mode = "shrink"
    size = start
    layout = estimate_text_layout(item, box, size)
    if fit_mode == "shrink":
        while size > min_size and (
            layout["overflow_height_ratio"] > args.max_overflow_ratio
            or layout["overflow_width_ratio"] > args.max_overflow_ratio
        ):
            size = max(min_size, size - 0.5)
            layout = estimate_text_layout(item, box, size)
    return round(size, 2), layout, role


def add_text(slide, item: dict[str, Any], box: tuple[float, float, float, float], default_font: str, args: argparse.Namespace) -> dict[str, Any]:
    x, y, w, h = box
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    effective_font, layout, role = choose_effective_font(item, box, args)
    vertical = str(item.get("vertical_anchor", "middle")).lower()
    shape.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(vertical, MSO_ANCHOR.MIDDLE)
    if item.get("rotation") is not None:
        shape.rotation = float(item.get("rotation"))
    tf = shape.text_frame
    tf.clear()
    role = infer_text_role(item)
    tf.word_wrap = bool(item.get("word_wrap", role not in NO_WRAP_DEFAULT_ROLES))
    ml, mr, mt, mb = text_margins_in(item, box)
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)
    lines = str(item.get("text", "")).split("\n")
    p = tf.paragraphs[0]
    p.text = lines[0] if lines else ""
    align = str(item.get("align", "left")).lower()
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)
    for line in lines[1:]:
        next_p = tf.add_paragraph()
        next_p.text = line
        next_p.alignment = p.alignment
    for para in tf.paragraphs:
        para.space_before = Pt(0)
        para.space_after = Pt(0)
        if item.get("line_spacing_pt") is not None:
            para.line_spacing = Pt(float(item.get("line_spacing_pt")))
        elif item.get("line_spacing") is not None:
            raw_spacing = float(item.get("line_spacing"))
            para.line_spacing = Pt(raw_spacing) if raw_spacing > 3 else raw_spacing
        for run in para.runs:
            run.font.name = item.get("font_face", default_font)
            run.font.size = Pt(effective_font)
            run.font.bold = bool(item.get("bold", False))
            run.font.italic = bool(item.get("italic", False))
            run.font.color.rgb = hex_to_rgb(item.get("color"), "003B7A")
    return {
        "id": item.get("id"),
        "slide": int(item.get("slide", 1)),
        "z_index": float(item.get("z_index", item.get("z", 0))),
        "text": str(item.get("text", "")),
        "text_role": role,
        "bbox_in": [round(v, 4) for v in box],
        "font_size": float(item.get("font_size", 14)),
        "effective_font_size": effective_font,
        "fit_mode": str(item.get("fit_mode", "shrink")).lower(),
        "min_font_size": role_min_font(role, item, args),
        "align": str(item.get("align", "left")).lower(),
        "vertical_anchor": vertical,
        "margins_in": [round(v, 4) for v in (ml, mr, mt, mb)],
        "parent_id": item.get("parent_id"),
        "text_overlay_allowed": bool(item.get("text_overlay_allowed", False)),
        **layout,
    }


def add_native_shape(slide, item: dict[str, Any], box: tuple[float, float, float, float]) -> None:
    x, y, w, h = box
    shape_name = str(item.get("shape", "round_rect")).lower()
    shape_type = {
        "rect": MSO_SHAPE.RECTANGLE,
        "rectangle": MSO_SHAPE.RECTANGLE,
        "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL,
        "oval": MSO_SHAPE.OVAL,
    }.get(shape_name, MSO_SHAPE.ROUNDED_RECTANGLE)
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if item.get("rotation") is not None:
        shape.rotation = float(item.get("rotation"))
    fill = str(item.get("fill", "FFFFFF"))
    if fill.lower() in {"none", "transparent"}:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill)
        transparency = float(item.get("fill_transparency", 0))
        if transparency:
            shape.fill.transparency = max(0, min(100, transparency)) / 100
    line = item.get("line", "C7DAF6")
    if str(line).lower() in {"none", "transparent"}:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_to_rgb(str(line), "C7DAF6")
        shape.line.width = Pt(float(item.get("line_width", 1)))
    if item.get("shadow"):
        try:
            shadow = shape.shadow
            shadow.inherit = False
            shadow.visible = True
            shadow.distance = Pt(float(item.get("shadow", {}).get("distance", 2)))
            shadow.blur_radius = Pt(float(item.get("shadow", {}).get("blur", 2)))
            shadow.transparency = float(item.get("shadow", {}).get("transparency", 55)) / 100
        except Exception:
            pass


def add_native_line(slide, item: dict[str, Any], box: tuple[float, float, float, float]) -> None:
    x, y, w, h = box
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x),
        Inches(y),
        Inches(x + w),
        Inches(y + h),
    )
    connector.line.color.rgb = hex_to_rgb(str(item.get("line", item.get("color", "5A93EA"))), "5A93EA")
    connector.line.width = Pt(float(item.get("line_width", 1.5)))


def fit_contain(img_path: Path, box: tuple[float, float, float, float]) -> tuple[float, float, float, float]:
    x, y, w, h = box
    with Image.open(img_path) as img:
        src_ratio = img.width / img.height
    slot_ratio = w / h
    if src_ratio > slot_ratio:
        fw = w
        fh = w / src_ratio
        fx = x
        fy = y + (h - fh) / 2
    else:
        fh = h
        fw = h * src_ratio
        fx = x + (w - fw) / 2
        fy = y
    return fx, fy, fw, fh


def add_asset(slide, item: dict[str, Any], asset: dict[str, Any], box: tuple[float, float, float, float], base_dir: Path) -> dict[str, Any]:
    source = asset.get("source_type")
    if source not in VALID_ASSET_SOURCES:
        raise ValueError(f"{item.get('id')} has invalid source_type: {source}")
    if asset.get("semantic_unit_count") != 1:
        raise ValueError(f"{item.get('id')} must have semantic_unit_count=1")
    img_path = rel_or_abs(base_dir, str(asset["asset_path"]))
    if not img_path.exists():
        raise FileNotFoundError(f"asset not found for {item.get('id')}: {img_path}")
    fx, fy, fw, fh = fit_contain(img_path, box)
    pic = slide.shapes.add_picture(str(img_path), Inches(fx), Inches(fy), width=Inches(fw), height=Inches(fh))
    if item.get("rotation") is not None:
        pic.rotation = float(item.get("rotation"))
    return {
        "semantic_unit_id": item.get("id"),
        "asset_path": str(img_path),
        "slot_in": [round(v, 4) for v in box],
        "fitted_in": [round(v, 4) for v in (fx, fy, fw, fh)],
    }


def get_items(inventory: dict[str, Any]) -> list[dict[str, Any]]:
    if isinstance(inventory.get("slides"), list):
        items = []
        for idx, slide in enumerate(inventory["slides"], start=1):
            for item in slide.get("items", []):
                item = dict(item)
                item.setdefault("slide", idx)
                items.append(item)
        return sorted(items, key=lambda item: (int(item.get("slide", 1)), float(item.get("z_index", item.get("z", 0)))))
    return sorted([dict(item) for item in inventory.get("items", [])], key=lambda item: (int(item.get("slide", 1)), float(item.get("z_index", item.get("z", 0)))))


def intersect_area(a: list[float], b: list[float]) -> float:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0.0, min(ay + ah, by + bh) - max(ay, by))
    return ix * iy


def area(box: list[float]) -> float:
    return max(0.0, box[2]) * max(0.0, box[3])


def qa_layout(
    records: list[dict[str, Any]],
    text_records: list[dict[str, Any]],
    args: argparse.Namespace,
    slide_w: float,
    slide_h: float,
) -> dict[str, Any]:
    errors: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    by_id = {rec["id"]: rec for rec in records if rec.get("id")}

    for rec in records:
        x, y, w, h = rec["bbox_in"]
        if x < -0.02 or y < -0.02 or x + w > slide_w + 0.02 or y + h > slide_h + 0.02:
            warnings.append({"type": "bounds", "id": rec.get("id"), "slide": rec.get("slide"), "bbox_in": rec["bbox_in"]})

    for text in text_records:
        overflow = max(float(text["overflow_height_ratio"]), float(text["overflow_width_ratio"]))
        if overflow > args.max_overflow_ratio:
            errors.append({
                "type": "text_overflow",
                "id": text["id"],
                "slide": text["slide"],
                "overflow_ratio": round(overflow, 4),
                "effective_font_size": text["effective_font_size"],
                "bbox_in": text["bbox_in"],
            })
        min_allowed = role_min_font(str(text["text_role"]), {"min_font_size": text.get("min_font_size")}, args)
        if text["text_role"] in {"title", "subtitle", "header"} and text["effective_font_size"] < min_allowed:
            errors.append({"type": "title_font_too_small", "id": text["id"], "slide": text["slide"], "effective_font_size": text["effective_font_size"]})
        elif text["text_role"] == "micro_label":
            pass
        elif text["text_role"] not in {"title", "subtitle", "header"} and text["effective_font_size"] < args.min_body_font:
            warnings.append({"type": "body_font_too_small", "id": text["id"], "slide": text["slide"], "effective_font_size": text["effective_font_size"]})

        parent_id = text.get("parent_id")
        if parent_id and parent_id in by_id:
            parent = by_id[parent_id]["bbox_in"]
            tx, ty, tw, th = text["bbox_in"]
            px, py, pw, ph = parent
            if tx < px - 0.03 or ty < py - 0.03 or tx + tw > px + pw + 0.03 or ty + th > py + ph + 0.03:
                warnings.append({"type": "parent_containment", "id": text["id"], "parent_id": parent_id, "slide": text["slide"]})

    for idx, a in enumerate(text_records):
        if not a.get("text", "").strip():
            continue
        for b in text_records[idx + 1:]:
            if a["slide"] != b["slide"] or not b.get("text", "").strip():
                continue
            if a.get("parent_id") and a.get("parent_id") == b.get("parent_id"):
                continue
            ia = intersect_area(a["estimated_text_bbox_in"], b["estimated_text_bbox_in"])
            min_area = max(0.0001, min(area(a["estimated_text_bbox_in"]), area(b["estimated_text_bbox_in"])))
            overlap = ia / min_area
            if overlap >= args.collision_threshold:
                warnings.append({
                    "type": "text_collision",
                    "slide": a["slide"],
                    "a": a["id"],
                    "b": b["id"],
                    "overlap_share": round(overlap, 4),
                })

    for text in text_records:
        for rec in records:
            if rec["slide"] != text["slide"] or rec["class"] not in {"imagegen_asset", "api_generated_asset", "provided_asset"}:
                continue
            role = str(rec.get("asset_role") or rec.get("role") or "").lower()
            if text.get("text_overlay_allowed") or rec.get("text_overlay_allowed"):
                continue
            if role in {"decor", "decoration", "underlay", "background", "icon", "badge"}:
                continue
            ia = intersect_area(text["estimated_text_bbox_in"], rec["bbox_in"])
            if ia / max(0.0001, area(text["estimated_text_bbox_in"])) >= 0.12:
                errors.append({
                    "type": "text_on_image",
                    "slide": text["slide"],
                    "text_id": text["id"],
                    "image_id": rec["id"],
                    "overlap_share": round(ia / max(0.0001, area(text["estimated_text_bbox_in"])), 4),
                })

    return {
        "status": "FAIL" if errors else "PASS_WITH_WARNINGS" if warnings else "PASS",
        "error_count": len(errors),
        "warning_count": len(warnings),
        "errors": errors,
        "warnings": warnings,
    }


def main() -> None:
    args = parse_args()
    inventory_path = Path(args.inventory)
    manifest_path = Path(args.manifest)
    base_dir = Path(args.base_dir) if args.base_dir else inventory_path.parent
    inventory = read_json(inventory_path)
    manifest = read_json(manifest_path)
    layout_rules = read_json(Path(args.layout_rules)) if args.layout_rules else {}
    manifest_by_id = {item["semantic_unit_id"]: item for item in manifest}
    slide_size_px = inventory.get("slide_size_px") or inventory.get("canvas_px") or [1920, 1080]
    default_font = (
        layout_rules.get("font_policy", {}).get("default_font_face")
        or inventory.get("font_face")
        or "Arial Unicode MS"
    )

    prs = Presentation()
    prs.slide_width = Inches(args.slide_width)
    prs.slide_height = Inches(args.slide_height)
    blank = prs.slide_layouts[6]
    slides = {}
    report = {
        "text_objects": 0,
        "native_layout_objects": 0,
        "asset_objects": 0,
        "text_placements": [],
        "native_layout_placements": [],
        "asset_placements": [],
    }
    records: list[dict[str, Any]] = []

    for item in get_items(inventory):
        slide_no = int(item.get("slide", 1))
        while len(slides) < slide_no:
            slides[len(slides) + 1] = prs.slides.add_slide(blank)
        slide = slides[slide_no]
        cls = item.get("class") or item.get("type")
        bbox = item.get("bbox_px") or item.get("bbox")
        if not bbox:
            raise ValueError(f"item {item.get('id')} missing bbox_px")
        box = px_box_to_inches(bbox, slide_size_px, args.slide_width, args.slide_height)
        if cls == "text":
            text_report = add_text(slide, item, box, default_font, args)
            report["text_placements"].append(text_report)
            report["text_objects"] += 1
        elif cls == "layout_native":
            add_native_shape(slide, item, box)
            report["native_layout_placements"].append({
                "id": item.get("id"),
                "slide": slide_no,
                "z_index": float(item.get("z_index", item.get("z", 0))),
                "bbox_in": [round(v, 4) for v in box],
                "shape": item.get("shape", "round_rect"),
            })
            report["native_layout_objects"] += 1
        elif cls in {"line_native", "connector_native"}:
            add_native_line(slide, item, box)
            report["native_layout_placements"].append({
                "id": item.get("id"),
                "slide": slide_no,
                "z_index": float(item.get("z_index", item.get("z", 0))),
                "bbox_in": [round(v, 4) for v in box],
                "shape": "line",
            })
            report["native_layout_objects"] += 1
        elif cls in {"imagegen_asset", "api_generated_asset", "provided_asset"}:
            asset = manifest_by_id.get(item.get("id"))
            if not asset:
                raise ValueError(f"no manifest entry for asset item: {item.get('id')}")
            report["asset_placements"].append(add_asset(slide, item, asset, box, base_dir))
            report["asset_objects"] += 1
        elif cls == "unresolved":
            raise ValueError(f"unresolved item cannot be built: {item.get('id')}")
        else:
            raise ValueError(f"unsupported item class for {item.get('id')}: {cls}")
        records.append({
            "id": item.get("id"),
            "slide": slide_no,
            "class": cls,
            "z_index": float(item.get("z_index", item.get("z", 0))),
            "bbox_in": [round(v, 4) for v in box],
            "parent_id": item.get("parent_id"),
            "role": item.get("role"),
            "asset_role": item.get("asset_role"),
            "text_overlay_allowed": bool(item.get("text_overlay_allowed", False)),
        })

    out_pptx = Path(args.out_pptx)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)
    report["pptx"] = str(out_pptx)
    report["slide_count"] = len(slides)
    report["layout_qa"] = qa_layout(records, report["text_placements"], args, args.slide_width, args.slide_height)
    report["layout_policy"] = {
        "default_font_face": default_font,
        "min_body_font": args.min_body_font,
        "min_title_font": args.min_title_font,
        "max_overflow_ratio": args.max_overflow_ratio,
        "collision_threshold": args.collision_threshold,
    }
    if args.report:
        report_path = Path(args.report)
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(report, ensure_ascii=False))
    if args.fail_on_layout_qa and report["layout_qa"]["error_count"]:
        raise SystemExit(2)


if __name__ == "__main__":
    main()
